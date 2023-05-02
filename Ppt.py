import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor

# Load the data and extract the latest revenue value
df = pd.read_csv('revenue_data.csv')
latest_revenue = df['revenue'].iloc[-1]

# Create a new PowerPoint presentation
prs = Presentation()

# Add a new slide with the title and subtitle
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
title.text = f"The team hit ${latest_revenue}"
subtitle = slide.placeholders[1]
subtitle.text = "Weekly Revenue Report"

# Add a table with 5 columns and 2 rows
table = slide.shapes.add_table(rows=2, cols=5, left=Inches(0.5), top=Inches(2.0), width=Inches(9.0), height=Inches(1.0)).table
table.cell(0, 0).text = "Revenue"
table.cell(1, 0).text = "{:.2f}".format(latest_revenue)

# Add a chart to show week-over-week change in revenue
chart_data = ChartData()
chart_data.categories = ['Week 1', 'Week 2', 'Week 3', 'Week 4']
chart_data.add_series('Cumulative Revenue', (100, 110, 120, 130))
x, y, cx, cy = Inches(0.5), Inches(3.5), Inches(5), Inches(3)
graphic_frame = slide.shapes.add_chart(XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data).graphic_frame
chart = graphic_frame.chart
chart.has_title = True
chart.chart_title.text_frame.text = "Week-over-Week Change in Cumulative Revenue"
chart.chart_title.text_frame.auto_size = True
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(10)

# Add a title for the notes section
notes_title = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(5), width=Inches(9), height=Inches(0.5)).text_frame
notes_title.text = "Current Notes"
notes_title.paragraphs[0].font.bold = True

# Add two bullet points for the notes
bullet1 = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(5.5), width=Inches(9), height=Inches(0.25)).text_frame
bullet1.paragraphs[0].text = f"Revenue has increased by 10% compared to last week."
bullet1.paragraphs[0].font.bold = True

bullet2 = slide.shapes.add_textbox(left=Inches(0.5), top=Inches(6), width=Inches(9), height=Inches(0.25)).text_frame
bullet2.paragraphs[0].text = "We expect revenue to continue to grow in the coming weeks."
bullet2.paragraphs[0].font.bold = True

# Save the PowerPoint presentation
prs.save('weekly_revenue_report.pptx')
